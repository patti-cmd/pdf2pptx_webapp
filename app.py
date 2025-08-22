import os
import tempfile
from flask import Flask, render_template, request, send_file, redirect, url_for, flash
from werkzeug.utils import secure_filename
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import io

ALLOWED_EXTENSIONS = {'.pdf'}
MAX_CONTENT_LENGTH = int(os.environ.get("MAX_CONTENT_LENGTH_BYTES", 50 * 1024 * 1024))
DEFAULT_DPI = int(os.environ.get("DEFAULT_DPI", 300))
POPPLER_PATH = os.environ.get("POPPLER_PATH")

app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get("FLASK_SECRET_KEY", "change-me")
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

def allowed_file(filename):
    _, ext = os.path.splitext(filename.lower())
    return ext in ALLOWED_EXTENSIONS

def pdf_to_pptx(pdf_path, dpi=DEFAULT_DPI, poppler_path=POPPLER_PATH):
    images = convert_from_path(pdf_path, dpi=dpi, poppler_path=poppler_path)
    if not images:
        raise RuntimeError("No pages found in PDF.")

    prs = Presentation()

    first = images[0]
    width_in = first.width / dpi
    height_in = first.height / dpi
    prs.slide_width = Inches(width_in)
    prs.slide_height = Inches(height_in)

    blank = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory() as tmpdir:
        for idx, img in enumerate(images):
            slide = prs.slides.add_slide(blank)
            img_path = os.path.join(tmpdir, f"page_{idx+1}.png")
            img.save(img_path, "PNG")
            slide.shapes.add_picture(
                img_path,
                left=Inches(0),
                top=Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
        out_path = os.path.join(tmpdir, "output.pptx")
        prs.save(out_path)
        with open(out_path, "rb") as f:
            return f.read()

@app.route("/", methods=["GET", "POST"])
def index():
    if request.method == "POST":
        file = request.files.get("pdf")
        dpi = request.form.get("dpi", type=int) or DEFAULT_DPI

        if not file or file.filename == "":
            flash("Please choose a PDF to upload.")
            return redirect(url_for("index"))

        filename = secure_filename(file.filename)
        if not allowed_file(filename):
            flash("Only .pdf files are allowed.")
            return redirect(url_for("index"))

        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                pdf_path = os.path.join(tmpdir, filename)
                file.save(pdf_path)

                pptx_bytes = pdf_to_pptx(pdf_path, dpi=dpi)
                stamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
                out_name = f"{os.path.splitext(filename)[0]}_{stamp}.pptx"
                return send_file(
                    io.BytesIO(pptx_bytes),
                    mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    as_attachment=True,
                    download_name=out_name
                )
        except Exception as e:
            flash(f"Conversion failed: {e}")
            return redirect(url_for("index"))

    return render_template("upload.html", default_dpi=DEFAULT_DPI, max_mb=app.config['MAX_CONTENT_LENGTH']//(1024*1024))

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)), debug=True)
